"""PPTX renderer — one slide per source page.

Blocks placed by `bbox` scaled to slide dimensions. Figures → picture shape
from inline crops; tables → native PowerPoint tables; text → text frames.
`reading_order` drives z-index (later blocks on top).
"""
from __future__ import annotations

from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from ocr_agentic_engine.renderers._ordering import iter_blocks
from ocr_agentic_engine.types import Annotations, DocumentRepresentation, OCRBlock

_SLIDE_W = 9144000   # 10in in EMU
_SLIDE_H = 6858000   # 7.5in in EMU


def render_pptx_slide(doc: DocumentRepresentation, ann: Annotations) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)
    blank = prs.slide_layouts[6]

    page_dims = {p.page: (p.width, p.height, p.path) for p in doc.page_images}
    slides_by_page: dict[int, any] = {}

    for idx, b in iter_blocks(doc, ann):
        if idx in ann.handwritten:
            continue
        if b.page not in page_dims:
            continue
        w, h, path = page_dims[b.page]
        if b.page not in slides_by_page:
            slides_by_page[b.page] = prs.slides.add_slide(blank)
        slide = slides_by_page[b.page]

        x0, y0, x1, y1 = b.bbox
        sx = _SLIDE_W / w; sy = _SLIDE_H / h
        left, top = Emu(int(x0 * sx)), Emu(int(y0 * sy))
        width, height = Emu(int((x1 - x0) * sx)), Emu(int((y1 - y0) * sy))

        if b.kind == "figure":
            _add_figure(slide, path, b, left, top, width, height)
        elif b.kind == "table" and b.table_cells:
            _add_table(slide, b, left, top, width, height)
        else:
            text = ann.corrections.get(idx, b.text)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tb.text_frame.text = text

    buf = BytesIO(); prs.save(buf); return buf.getvalue()


def _add_figure(slide, page_path, b: OCRBlock, left, top, width, height):
    try:
        with Image.open(page_path) as im:
            x0, y0, x1, y1 = b.bbox
            crop = im.crop((max(0, x0), max(0, y0), min(im.width, x1), min(im.height, y1)))
            buf = BytesIO(); crop.save(buf, format="PNG"); buf.seek(0)
        slide.shapes.add_picture(buf, left, top, width=width, height=height)
    except Exception:
        slide.shapes.add_textbox(left, top, width, height).text_frame.text = "[figure]"


def _add_table(slide, b: OCRBlock, left, top, width, height):
    rows: dict[int, dict[int, str]] = {}
    for c in b.table_cells or []:
        rows.setdefault(c.row, {})[c.col] = c.text
    if not rows:
        slide.shapes.add_textbox(left, top, width, height).text_frame.text = b.text
        return
    n_rows = max(rows) + 1
    n_cols = max((max(r) for r in rows.values()), default=0) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for r, cols in rows.items():
        for c, text in cols.items():
            tbl.cell(r, c).text = text
