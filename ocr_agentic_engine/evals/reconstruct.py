"""Reconstruction eval — exercise every renderer and score fidelity.

Answers: "does the L3 renderer preserve what L1 gave it?"

For each scan_conversion fixture:
  1. Run L1 (OCR) + heuristic L2 (heading / list annotators only — no LLM).
  2. Render to every format in `REGISTRY` (pdf, docx, txt, html, xlsx, pptx).
  3. Extract text back from each output file (pdf→pdfium, docx→python-docx,
     html→strip tags, txt→utf-8, xlsx→cells, pptx→shape text).
  4. Score text equality vs `doc.raw_text`. 1.0 = lossless round-trip.

Also records per-format render time and output bytes.

This runs **without an LLM** — it's a pure L1+L3 test. For the L2 agent path
(classifier / icr_corrector / handwriting_filter), use `evals run`.

Usage:
    uv run python -m ocr_agentic_engine.evals reconstruct
    uv run python -m ocr_agentic_engine.evals reconstruct --fixture Invoice1.png
"""
from __future__ import annotations

import argparse
import json
import re
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from ocr_agentic_engine.agents.annotators import HeadingClassifierAgent, ListDetectorAgent
from ocr_agentic_engine.evals import metrics
from ocr_agentic_engine.evals.runner import FIXTURES_DIR, RESULTS_ROOT
from ocr_agentic_engine.renderers import REGISTRY, render
from ocr_agentic_engine.types import Annotations, DocumentRepresentation


@dataclass
class ReconstructCase:
    feature: str
    fixture: Path


def _discover() -> list[ReconstructCase]:
    out: list[ReconstructCase] = []
    for feature in ("scan_conversion", "classification", "handwriting_removal"):
        d = FIXTURES_DIR / feature
        if not d.is_dir():
            continue
        for p in sorted(d.iterdir()):
            if (p.is_file()
                    and p.suffix.lower() in {".pdf", ".png", ".jpg", ".jpeg"}
                    and ".annotated" not in p.suffixes
                    and not p.name.endswith(".golden.json")):
                out.append(ReconstructCase(feature, p))
    return out


def _ocr_stage():
    from paddleocr import PPStructureV3
    from ocr_agentic_engine.ocr_stage import OCRStage
    structure = PPStructureV3(use_doc_orientation_classify=False, use_doc_unwarping=False,
                               use_textline_orientation=False, use_seal_recognition=False,
                               use_formula_recognition=False, use_chart_recognition=False)
    return OCRStage(structure=structure)


# --- per-format text extractors -------------------------------------------

def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_html(data: bytes) -> str:
    src = data.decode("utf-8", errors="replace")
    # drop <style> and <script> blocks (their text content would otherwise pollute)
    src = re.sub(r"<(style|script)\b[^>]*>.*?</\1>", " ", src, flags=re.DOTALL | re.IGNORECASE)
    no_tags = re.sub(r"<[^>]+>", " ", src)
    return re.sub(r"\s+", " ", no_tags).strip()


def _extract_docx(data: bytes) -> str:
    from docx import Document
    d = Document(BytesIO(data))
    parts: list[str] = []
    for p in d.paragraphs:
        parts.append(p.text)
    for tbl in d.tables:
        for row in tbl.rows:
            parts.extend(c.text for c in row.cells)
    return "\n".join(x for x in parts if x)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(BytesIO(data), data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.extend(c.text for c in row.cells)
    return "\n".join(x for x in parts if x)


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        r = PdfReader(BytesIO(data))
        return "\n".join(p.extract_text() or "" for p in r.pages)
    except Exception:
        return ""


EXTRACTORS = {
    "txt": _extract_txt, "html": _extract_html, "docx": _extract_docx,
    "xlsx": _extract_xlsx, "pptx": _extract_pptx, "pdf": _extract_pdf,
}


# --- main loop ------------------------------------------------------------

def run(fixture_filter: str | None = None) -> Path:
    stage = _ocr_stage()
    heading = HeadingClassifierAgent(); lists = ListDetectorAgent()

    started = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    out = RESULTS_ROOT / f"reconstruct-{started}"
    out.mkdir(parents=True, exist_ok=True)
    report_path = out / "report.jsonl"

    with report_path.open("w") as fp:
        for case in _discover():
            if fixture_filter and fixture_filter not in case.fixture.name:
                continue

            tmp_root = Path("/tmp/engine-tmp") / f"recon-{case.fixture.stem}"
            tmp_root.mkdir(parents=True, exist_ok=True)

            print(f"\n[reconstruct] {case.feature}/{case.fixture.name}")
            t0 = time.perf_counter()
            try:
                doc = stage.run(case.fixture, tmp_root)
            except Exception as e:
                fp.write(json.dumps({"fixture": case.fixture.name, "ocr_error": str(e)}) + "\n")
                print(f"  OCR FAILED: {e}")
                continue
            ocr_ms = (time.perf_counter() - t0) * 1000
            ann = Annotations(headings=heading.annotate(doc), lists=lists.annotate(doc))

            reference = _norm(doc.raw_text)
            per_format: dict[str, dict] = {}
            for fmt in REGISTRY:
                tr = time.perf_counter()
                try:
                    data, _ = render(fmt, doc, ann)
                    render_ms = (time.perf_counter() - tr) * 1000
                    extracted = _norm(EXTRACTORS[fmt](data))
                    score = metrics.text_tree_equality(extracted, reference)["score"]
                    per_format[fmt] = {
                        "bytes": len(data),
                        "render_ms": round(render_ms, 1),
                        "roundtrip_score": round(score, 3),
                    }
                    print(f"  {fmt:5s}: score={score:.3f} bytes={len(data):>7d} ms={render_ms:.0f}")
                except Exception as e:
                    per_format[fmt] = {"error": str(e)}
                    print(f"  {fmt:5s}: ERROR {e}")

            fp.write(json.dumps({
                "feature": case.feature,
                "fixture": case.fixture.name,
                "blocks": len(doc.blocks),
                "pages": doc.page_count,
                "ocr_ms": round(ocr_ms, 1),
                "raw_text_chars": len(doc.raw_text),
                "formats": per_format,
            }) + "\n")

    (out / "summary.md").write_text(_summarise(report_path))
    print(f"\n[reconstruct] report → {out}")
    return out


def _norm(s: str) -> str:
    return " ".join((s or "").split()).lower()


def _summarise(report_path: Path) -> str:
    rows = [json.loads(l) for l in report_path.read_text().splitlines()]
    lines = [f"# Reconstruction report — {report_path.parent.name}", ""]
    lines.append("| fixture | blocks | format | roundtrip | bytes | render_ms |")
    lines.append("|---|---:|---|---:|---:|---:|")
    for row in rows:
        if "formats" not in row:
            lines.append(f"| {row.get('fixture','?')} | — | OCR | FAIL | — | — |")
            continue
        for fmt, v in row["formats"].items():
            if "error" in v:
                lines.append(f"| {row['fixture']} | {row['blocks']} | {fmt} | ERROR | — | — |")
            else:
                lines.append(f"| {row['fixture']} | {row['blocks']} | {fmt} "
                              f"| {v['roundtrip_score']:.3f} | {v['bytes']} | {v['render_ms']} |")
    # per-format average
    per_fmt: dict[str, list[float]] = {}
    for row in rows:
        for fmt, v in row.get("formats", {}).items():
            if "roundtrip_score" in v:
                per_fmt.setdefault(fmt, []).append(v["roundtrip_score"])
    lines.append("\n## Per-format averages\n")
    lines.append("| format | mean roundtrip | n |")
    lines.append("|---|---:|---:|")
    for fmt, scores in sorted(per_fmt.items()):
        lines.append(f"| {fmt} | {sum(scores)/len(scores):.3f} | {len(scores)} |")
    return "\n".join(lines) + "\n"


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--fixture", help="substring filter on fixture filename")
    args = p.parse_args(argv)
    run(args.fixture)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
